"""
StudHub Demo Deck Generator
Run: pip install python-pptx && python paper/studhub_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Brand colours ──────────────────────────────────────────────────────────
RED_PRIMARY   = RGBColor(0xE6, 0x39, 0x46)   # StudHub red
DARK_TEXT     = RGBColor(0x21, 0x21, 0x21)
GRAY_TEXT     = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_BG      = RGBColor(0xF8, 0xF9, 0xFA)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, text: str, bg=RED_PRIMARY, font_size=40, top=True):
    """Full-width header bar."""
    h = Inches(1.0)
    top_offset = Inches(0) if top else SLIDE_H - h
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), top_offset,
        SLIDE_W, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


def add_body_text(slide, lines: list[str], left=Inches(1), top=Inches(1.8), width=Inches(11), height=Inches(5)):
    """Add a multi-line body text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(20)
        run.font.color.rgb = DARK_TEXT
        run.font.name = "Calibri"


def add_bullet_list(slide, items: list[str], left=Inches(1.2), top=Inches(1.8), width=Inches(10.5)):
    """Add bullet points."""
    txb = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(19)
        run.font.color.rgb = DARK_TEXT
        run.font.name = "Calibri"


def add_footer(slide, text="StudHub — SEAIT — 2026"):
    """Slide number / footer bar."""
    h = Inches(0.35)
    shape = slide.shapes.add_shape(1, Inches(0), SLIDE_H - h, SLIDE_W, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY_TEXT
    run.font.name = "Calibri"


def add_two_column(slide, left_title, left_items, right_title, right_items,
                   left=Inches(0.8), top=Inches(1.8), col_w=Inches(5.8)):
    """Two-column layout."""
    # Left column
    lb = slide.shapes.add_textbox(left, top, col_w, Inches(0.5)).text_frame
    lp = lb.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_title
    lr.font.size = Pt(22)
    lr.font.bold = True
    lr.font.color.rgb = RED_PRIMARY
    lr.font.name = "Calibri"

    add_bullet_list(slide, left_items,
                    left=left, top=top + Inches(0.55), width=col_w - Inches(0.2))

    # Right column
    rb = slide.shapes.add_textbox(left + col_w + Inches(0.6), top, col_w, Inches(0.5)).text_frame
    rp = rb.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_title
    rr.font.size = Pt(22)
    rr.font.bold = True
    rr.font.color.rgb = RED_PRIMARY
    rr.font.name = "Calibri"

    add_bullet_list(slide, right_items,
                    left=left + col_w + Inches(0.6), top=top + Inches(0.55), width=col_w - Inches(0.2))


def new_slide(prs, title="", layout_idx=6, notes=""):
    """Add a blank slide with optional speaker notes.

    If 'title' is passed and no explicit notes are given,
    the title is used as the speaker note.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    set_bg(slide, LIGHT_BG)
    note_text = notes if notes else title
    if note_text:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = note_text
    return slide


def build_deck(output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ────────────────────────────────────────────────────
    s = new_slide(prs, "Title")
    set_bg(s, RED_PRIMARY)

    # White text centred
    txb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate([
        ("StudHub", 72, True, WHITE),
        ("Cross-Program Academic Resource Exchange for SEAIT", 28, False, RGBColor(0xFF, 0xCC, 0xCC)),
        ("", 14, False, WHITE),
        ("[Student Name]  |  South East Asian Institute of Technology  |  May 2026", 16, False, WHITE),
    ]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    # ── Slide 2: Problem Statement ─────────────────────────────────────────
    s = new_slide(prs, "The Problem — studhub: students can't find resources across programs")
    add_header_bar(s, "The Problem")
    add_body_text(s, [
        "When a BSCE student needs a Data Structures reviewer that a BSIT",
        "student already has, the current workflow is:",
        "",
        "  1.  Post in a random Facebook group",
        "  2.  Scroll Messenger chat history",
        "  3.  Give up and re-make the reviewer from scratch",
        "",
        "No directory of who has what.  No cross-program discovery.",
        "  No signal for the institution.",
    ])
    add_footer(s)

    # ── Slide 3: Solution ──────────────────────────────────────────────────
    s = new_slide(prs, "StudHub Solution Overview — Chat + Catalog + Routing Engine")
    add_header_bar(s, "StudHub — Solution Overview")
    add_two_column(
        s,
        "Per-Program Chat",
        [
            "Real-time rooms per program & year",
            "@display_name mentions + notifications",
            "File attachments (25 MB, MIME-validated)",
            "Suspended users blocked at HTTP + WS level",
        ],
        "Resource Catalog",
        [
            "Subject-tagged entries across all programs",
            "Full-text search + shelves",
            "PDF upload with per-user watermarking",
            "Program-only or school-wide visibility",
        ],
    )
    add_footer(s)

    # Second half of slide 3
    s = new_slide(prs, "StudHub Solution Overview cont. — Reputation + Moderation")
    add_header_bar(s, "StudHub — Solution Overview (cont.)")
    add_two_column(
        s,
        "Request Routing Engine",
        [
            "Cross-program request matching",
            "Weighted scoring: curriculum + resources + history",
            "historicalFulfillmentRate() feedback loop",
            "Auto-notifies top users in qualifying programs",
        ],
        "Reputation + Moderation",
        [
            "Karma: upload +5, fulfill +2, report confirmed -5",
            "Badge tiers: Bronze / Silver / Gold",
            "Report system across messages, resources, users",
            "Audit log with message snapshot before hide",
        ],
    )
    add_footer(s)

    # ── Slide 4: Architecture ───────────────────────────────────────────────
    s = new_slide(prs, "Architecture — Laravel 11 + Reverb + MySQL + Redis")
    add_header_bar(s, "Architecture")
    add_body_text(s, [
        "Laravel 11 + Livewire 3 + Alpine.js + Tailwind CSS",
        "Laravel Reverb (WebSocket) for real-time messaging",
        "MySQL 8 (single-tenant) | Redis (queues + pub/sub)",
        "PDF watermarking via FPDI (pure PHP, no Imagick required)",
        "Pest 3 testing | PHPStan Level 6 | Pint PSR-12",
        "GitHub Actions CI: Pint → PHPStan → Pest",
    ])
    add_footer(s)

    # ── Slide 5: Per-Program Chat ──────────────────────────────────────────
    s = new_slide(prs, "Feature 1 — Per-Program Chat: real-time rooms, @mentions, file attachments")
    add_header_bar(s, "Key Feature #1 — Per-Program Chat")
    add_two_column(
        s,
        "Room Structure",
        [
            "Auto-provisioned on program seed",
            "Per-program rooms + per-year sub-channels",
            "Program-scope enforced at HTTP + WebSocket",
            "Year-level access control (Y3 can't see Y1 room)",
        ],
        "@Mentions + Notifications",
        [
            "@display_name parsed in message body",
            "Email notification sent immediately",
            "Deduplicated mentions (one notification per user)",
            "MIME-validated file attachments (25 MB limit)",
        ],
    )
    add_footer(s)

    # ── Slide 6: Resource Catalog ─────────────────────────────────────────
    s = new_slide(prs, "Feature 2 — Resource Catalog: subject-tagged, searchable, watermarked PDFs")
    add_header_bar(s, "Key Feature #2 — Resource Catalog")
    add_two_column(
        s,
        "Upload + Discovery",
        [
            "Subject-tagged (not program-tagged)",
            "Course code + year taken metadata",
            "MySQL FULLTEXT search across title + description",
            "Filter by subject, type, program, year level",
        ],
        "Shelves + PDF Watermarking",
        [
            "Save/bookmark resources to personal shelf",
            "Per-user PDF watermark at download time",
            "Watermark: 'Downloaded by {name} on {date}'",
            "Thumbnail preview for PDF resources",
        ],
    )
    add_footer(s)

    # ── Slide 7: Routing Engine ─────────────────────────────────────────────
    s = new_slide(prs, "Feature 3 — Request Routing Engine: weighted scoring across 5 factors")
    add_header_bar(s, "Key Feature #3 — Request Routing Engine")
    add_body_text(s, [
        "Score per program:  score(p,s) = w_edge×weight + w_resource×resCount",
        "                           + w_history×fulfillmentRate + w_proximity×yearBonus",
        "                           + w_urgency×urgencyMult − penalty_self",
        "",
        "Default weights: w_edge=0.40  |  w_resource=0.25  |  w_history=0.20",
        "                 w_proximity=0.10  |  w_urgency=0.05  |  penalty=0.05",
        "",
        "PROGRAM_THRESHOLD = 0.35  (below which programs are not notified)",
        "CHAT_THRESHOLD = 0.65     (above which rooms get a cross-post)",
    ])
    add_footer(s)

    # ── Slide 8: Key Innovation ─────────────────────────────────────────────
    s = new_slide(prs, "Key Innovation — historicalFulfillmentRate: self-improving routing feedback loop")
    add_header_bar(s, "Key Innovation — historicalFulfillmentRate()")
    add_body_text(s, [
        "Fulfillment rate per program-subject pair, updated on every accepted offer:",
        "",
        "  accepted_count(program, subject)  ÷  total_offers(program, subject)",
        "",
        "New programs with no history get neutral rate of 0.5.",
        "  Programs with high acceptance rate get boosted in future routing.",
        "",
        "This creates a feedback loop: the more a program fulfills requests,",
        "the more routing weight it receives, driving more fulfilling users to help.",
    ])
    add_footer(s)

    # ── Slide 9: Moderation + Reputation ───────────────────────────────────
    s = new_slide(prs, "Moderation — School-scoped reports, karma system, audit log with message snapshot")
    add_header_bar(s, "Moderation & Reputation")
    add_two_column(
        s,
        "Moderation System",
        [
            "School-scoped Report global scope (F17)",
            "Reports for messages, resources, users",
            "Moderator program dashboard (SQL-filtered)",
            "Admin dashboard with sign-up + activity metrics",
            "Suspended users blocked from HTTP + WebSocket",
        ],
        "Karma + Badges",
        [
            "Upload: +5  |  Accept offer: +3  |  Fulfill: +2",
            "Report confirmed: -5  |  Resource archived: -3",
            "Bronze badge: 25 karma  |  Silver: 75  |  Gold: 150",
            "Leaderboard by program",
            "Atomic karma increment (F9 — no race conditions)",
        ],
    )
    add_footer(s)

    # ── Slide 10: Pilot Results ─────────────────────────────────────────────
    s = new_slide(prs, "Pilot Results — BSIT x BSCE x BSBA-MM, one semester at SEAIT")
    add_header_bar(s, "Pilot Results — [Semester TBD]")
    add_body_text(s, [
        "[To be populated after pilot completion — placeholders below]",
        "",
        "  Target programs:  BSIT  ×  BSCE  ×  BSBA-MM",
        "  Confirmed invitees:  ~30 students",
        "  Resources uploaded target:  ≥ 100",
        "  Cross-program fulfillments target:  ≥ 50",
        "  Median request → match time target:  < 24 hours",
        "",
        "Pilot to run one full semester at SEAIT, Tupi, South Cotabato.",
    ])
    add_footer(s)

    # ── Slide 11: Lessons Learned ─────────────────────────────────────────
    s = new_slide(prs, "Lessons Learned — what worked, what to change in v2")
    add_header_bar(s, "Lessons Learned")
    add_body_text(s, [
        "[To be populated after pilot — sample structure below]",
        "",
        "What worked:",
        "  • Per-program chat rooms created natural engagement",
        "  • Subject tagging made cross-program discovery feel natural",
        "  • Karma incentives drove early uploads from power users",
        "",
        "What we would change:",
        "  • Add resource version history earlier",
        "  • Invest more in onboarding UX before pilot",
        "  • Seed more sample resources before first student invite",
    ])
    add_footer(s)

    # ── Slide 12: Q&A ──────────────────────────────────────────────────────
    s = new_slide(prs)
    set_bg(s, RED_PRIMARY)
    txb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate([
        ("Questions & Answers", 56, True, WHITE),
        ("", 20, False, WHITE),
        ("Thank you for your time.", 26, False, RGBColor(0xFF, 0xCC, 0xCC)),
    ]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    add_footer(s)

    prs.save(output_path)
    print(f"[OK] Saved: {output_path}")


if __name__ == "__main__":
    out = os.path.join(os.path.dirname(__file__), "studhub_demo_deck.pptx")
    build_deck(out)